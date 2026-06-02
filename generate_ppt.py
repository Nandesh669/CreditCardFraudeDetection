from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Theme Colors (Premium Slate Dark Theme)
BG_COLOR = RGBColor(15, 23, 42)        # Slate 900 (#0F172A)
CARD_BG = RGBColor(30, 41, 59)         # Slate 800 (#1E293B)
TEXT_PRIMARY = RGBColor(248, 250, 252) # Slate 50 (#F8FAFC)
TEXT_SECONDARY = RGBColor(148, 163, 184) # Slate 400 (#94A3B8)
ACCENT_BLUE = RGBColor(14, 165, 233)   # Sky 500 (#0EA5E9)
ACCENT_RED = RGBColor(244, 63, 94)     # Rose 500 (#F43F5E)
ACCENT_GREEN = RGBColor(16, 185, 129)  # Emerald 500 (#10B981)
CARD_BORDER = RGBColor(71, 85, 105)    # Slate 600 (#475569)

def add_slide_header(slide, title_text, category_text="CREDIT CARD FRAUD DETECTION"):
    # Add a thin colored top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ACCENT_BLUE
    top_bar.line.color.rgb = ACCENT_BLUE
    
    # Category tracker at the top
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    cat_tf = cat_box.text_frame
    cat_tf.word_wrap = True
    cat_tf.margin_left = Inches(0)
    cat_tf.margin_top = Inches(0)
    p_cat = cat_tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = 'Segoe UI'
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_BLUE
    
    # Main Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.margin_left = Inches(0)
    title_tf.margin_top = Inches(0)
    p_title = title_tf.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_PRIMARY

def add_card(slide, left, top, width, height, title, body_text_list, accent_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    if accent_color:
        card.line.color.rgb = accent_color
        card.line.width = Pt(2.5)
    else:
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.2)
        
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.25)
    tf.margin_bottom = Inches(0.25)
    
    # Title paragraph
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    if accent_color:
        p_title.font.color.rgb = accent_color
    else:
        p_title.font.color.rgb = TEXT_PRIMARY
    p_title.space_after = Pt(12)
    p_title.alignment = PP_ALIGN.LEFT
    
    # Body paragraphs
    for text in body_text_list:
        p = tf.add_paragraph()
        p.text = text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_SECONDARY
        p.space_after = Pt(8)
        p.alignment = PP_ALIGN.LEFT

def create_title_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    left = Inches(1.5)
    top = Inches(1.8)
    width = Inches(10.333)
    height = Inches(3.9)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ACCENT_BLUE
    card.line.width = Pt(3)
    
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.5)
    tf.margin_bottom = Inches(0.5)
    
    p_title = tf.paragraphs[0]
    p_title.text = "AI-POWERED CREDIT CARD FRAUD DETECTION"
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_PRIMARY
    p_title.space_after = Pt(20)
    p_title.alignment = PP_ALIGN.CENTER
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Real-Time Transaction Monitoring & Machine Learning Anomaly Detection"
    p_sub.font.name = 'Segoe UI'
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = TEXT_SECONDARY
    p_sub.space_after = Pt(30)
    p_sub.alignment = PP_ALIGN.CENTER
    
    p_badge = tf.add_paragraph()
    p_badge.text = "POWERED BY MACHINE LEARNING & FLASK"
    p_badge.font.name = 'Segoe UI'
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = ACCENT_BLUE
    p_badge.alignment = PP_ALIGN.CENTER

def create_problem_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    add_slide_header(slide, "The Problem: Rising Financial Threat", "UNDERSTANDING THE THREAT")
    
    card_width = Inches(3.6)
    card_height = Inches(4.5)
    top = Inches(1.8)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    
    add_card(slide, 
             start_left, 
             top, 
             card_width, 
             card_height, 
             "Financial Impact", 
             [
                 "Billions of dollars are lost globally every year due to unauthorized card transactions.",
                 "Chargeback fees, operational overhead, and recovery costs multiply the direct losses for financial institutions."
             ], 
             accent_color=ACCENT_RED)
             
    add_card(slide, 
             start_left + card_width + gap, 
             top, 
             card_width, 
             card_height, 
             "Customer Trust", 
             [
                 "Security breaches damage the reputation of banks and merchant platforms.",
                 "Fraud events trigger user anxiety, leading to customer churn and decreased brand loyalty."
             ], 
             accent_color=ACCENT_RED)
             
    add_card(slide, 
             start_left + (card_width + gap) * 2, 
             top, 
             card_width, 
             card_height, 
             "Traditional Limits", 
             [
                 "Legacy rule-based systems are too slow and fail to catch modern, sophisticated patterns.",
                 "High rate of false positives blocks legitimate transactions, frustrating users."
             ], 
             accent_color=ACCENT_RED)

def create_solution_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    add_slide_header(slide, "The Solution: Real-Time Machine Learning", "AI-DRIVEN INTELLIGENCE")
    
    card_width = Inches(3.6)
    card_height = Inches(4.5)
    top = Inches(1.8)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    
    add_card(slide, 
             start_left, 
             top, 
             card_width, 
             card_height, 
             "Predictive Approach", 
             [
                 "Utilizing Machine Learning to analyze multidimensional transaction patterns in real-time.",
                 "Extracting features such as transaction amount, location anomalies, device signatures, and user account age."
             ], 
             accent_color=ACCENT_GREEN)
             
    add_card(slide, 
             start_left + card_width + gap, 
             top, 
             card_width, 
             card_height, 
             "Core Objective", 
             [
                 "Accurately classify transactions into Legitimate or Fraudulent states.",
                 "Maximize detection sensitivity (recall) while minimizing false alarms (precision optimization)."
             ], 
             accent_color=ACCENT_GREEN)
             
    add_card(slide, 
             start_left + (card_width + gap) * 2, 
             top, 
             card_width, 
             card_height, 
             "Strategic Benefits", 
             [
                 "⚡ Sub-second execution and detection times.",
                 "📉 Dramatic reduction in false positives.",
                 "🔄 Continuous learning and adaptation as fraud patterns evolve."
             ], 
             accent_color=ACCENT_GREEN)

def create_tech_stack_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    add_slide_header(slide, "Tools & Technologies Used", "SYSTEM TECH STACK")
    
    row1_width = Inches(3.6)
    row1_height = Inches(2.1)
    row1_top = Inches(1.8)
    row1_gap = Inches(0.4)
    row1_start = Inches(0.8)
    
    add_card(slide, row1_start, row1_top, row1_width, row1_height, "Machine Learning", ["scikit-learn", "Random Forest Classifier"], accent_color=ACCENT_BLUE)
    add_card(slide, row1_start + row1_width + row1_gap, row1_top, row1_width, row1_height, "Backend Framework", ["Flask (Python)", "joblib model hosting"], accent_color=ACCENT_BLUE)
    add_card(slide, row1_start + (row1_width + row1_gap) * 2, row1_top, row1_width, row1_height, "Frontend Client", ["HTML5, Vanilla CSS", "Glassmorphism aesthetics"], accent_color=ACCENT_BLUE)
    
    row2_width = Inches(5.6)
    row2_height = Inches(2.1)
    row2_top = Inches(4.3)
    row2_gap = Inches(0.5)
    row2_start = Inches(0.8)
    
    add_card(slide, row2_start, row2_top, row2_width, row2_height, "Data Processing", ["Pandas & NumPy", "CSV dataset parser and model feature scaling"], accent_color=ACCENT_BLUE)
    add_card(slide, row2_start + row2_width + row2_gap, row2_top, row2_width, row2_height, "Analytics & Visuals", ["Chart.js integration", "Dynamic prediction tracking"], accent_color=ACCENT_BLUE)

def create_pipeline_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    add_slide_header(slide, "The Machine Learning Pipeline", "MODEL LIFECYCLE")
    
    card_width = Inches(2.65)
    card_height = Inches(4.5)
    top = Inches(1.8)
    gap = Inches(0.38)
    start_left = Inches(0.8)
    
    add_card(slide, start_left, top, card_width, card_height, "01. Data Preprocessing", [
        "Read transactions from transactions.csv.",
        "Select features: amount, time, location, device, type, account age.",
        "Split dataset 80/20 for training/testing."
    ], accent_color=ACCENT_BLUE)
    
    add_card(slide, start_left + card_width + gap, top, card_width, card_height, "02. Model Training", [
        "Instantiate RandomForestClassifier.",
        "Fit model to training dataset to learn feature relationships and decision boundaries."
    ], accent_color=ACCENT_BLUE)
    
    add_card(slide, start_left + (card_width + gap) * 2, top, card_width, card_height, "03. Evaluation & Save", [
        "Validate model accuracy on the test set.",
        "Export the trained model to fraud_model.pkl using joblib for production use."
    ], accent_color=ACCENT_BLUE)
    
    add_card(slide, start_left + (card_width + gap) * 3, top, card_width, card_height, "04. Web Integration", [
        "Flask app loads fraud_model.pkl on startup.",
        "Predicts class in real-time when new data is posted via endpoint."
    ], accent_color=ACCENT_BLUE)

def create_architecture_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    add_slide_header(slide, "System Architecture & Data Flow", "UNDER THE HOOD")
    
    card_width = Inches(2.65)
    card_height = Inches(4.5)
    top = Inches(1.8)
    gap = Inches(0.38)
    start_left = Inches(0.8)
    
    add_card(slide, start_left, top, card_width, card_height, "1. User Dashboard", [
        "Enter transaction inputs (amount, time, account age, location, device, type) on the dashboard web interface."
    ], accent_color=ACCENT_BLUE)
    
    add_card(slide, start_left + card_width + gap, top, card_width, card_height, "2. Flask Backend", [
        "Receive the transaction data POST request.",
        "Convert web form parameters into a numpy array matching the model schema."
    ], accent_color=ACCENT_BLUE)
    
    add_card(slide, start_left + (card_width + gap) * 2, top, card_width, card_height, "3. Prediction Engine", [
        "Execute RandomForestClassifier prediction.",
        "Analyze features and yield discrete class prediction (0 = Safe, 1 = Fraud)."
    ], accent_color=ACCENT_BLUE)
    
    add_card(slide, start_left + (card_width + gap) * 3, top, card_width, card_height, "4. Real-time Response", [
        "Update the web template dynamically.",
        "Trigger warnings (Red/Green indicators) and refresh analytics charts."
    ], accent_color=ACCENT_BLUE)

def create_ui_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    add_slide_header(slide, "User Interface: Modern Dashboard", "FRONTEND DESIGN")
    
    card_width = Inches(3.6)
    card_height = Inches(4.5)
    top = Inches(1.8)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    
    add_card(slide, start_left, top, card_width, card_height, "Premium Design", [
        "Sleek dark mode theme designed for security operations.",
        "Glassmorphism UI elements featuring elegant borders, subtle shadows, and transparent panel layers."
    ], accent_color=ACCENT_BLUE)
    
    add_card(slide, start_left + card_width + gap, top, card_width, card_height, "Interactive Simulation", [
        "Fully interactive forms allow immediate evaluation of new transactions.",
        "Vibrant response styling: bright green highlights for verified safe transactions, glowing red alert bars for detected fraud."
    ], accent_color=ACCENT_BLUE)
    
    add_card(slide, start_left + (card_width + gap) * 2, top, card_width, card_height, "Data Visualization", [
        "Chart.js charts demonstrate system accuracy, prediction trends, and statistics.",
        "Clear visualizations make detection outcomes easily readable for risk analysts."
    ], accent_color=ACCENT_BLUE)

def create_features_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    add_slide_header(slide, "Key Features & Capabilities", "SYSTEM HIGHLIGHTS")
    
    card_width = Inches(3.6)
    card_height = Inches(4.5)
    top = Inches(1.8)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    
    add_card(slide, start_left, top, card_width, card_height, "Millisecond Latency", [
        "The model classifies transactions in under 5 milliseconds.",
        "Optimized Flask pipeline ensures immediate feedback for high-volume endpoints."
    ], accent_color=ACCENT_GREEN)
    
    add_card(slide, start_left + card_width + gap, top, card_width, card_height, "Robust Classification", [
        "Random Forest Classifier provides excellent performance on structured transaction tables.",
        "Handles non-linear feature relationships and avoids overfitting through ensemble bagging."
    ], accent_color=ACCENT_GREEN)
    
    add_card(slide, start_left + (card_width + gap) * 2, top, card_width, card_height, "Seamless Integration", [
        "Lightweight architecture is easily containerized (Docker-ready).",
        "Clean boundaries allow integration with live credit card processing streams."
    ], accent_color=ACCENT_GREEN)

def create_future_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    add_slide_header(slide, "Future Enhancements & Roadmap", "FUTURE OUTLOOK")
    
    card_width = Inches(3.6)
    card_height = Inches(4.5)
    top = Inches(1.8)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    
    add_card(slide, start_left, top, card_width, card_height, "1. Database Integration", [
        "Connect the system to persistent storage (PostgreSQL/MongoDB).",
        "Log all transactions, predictions, and model outcomes for continuous auditing and visualization."
    ], accent_color=ACCENT_BLUE)
    
    add_card(slide, start_left + card_width + gap, top, card_width, card_height, "2. Advanced ML Models", [
        "Incorporate advanced gradient boosting models (XGBoost, LightGBM) or deep neural networks.",
        "Train models on large-scale datasets with millions of rows for enhanced precision."
    ], accent_color=ACCENT_BLUE)
    
    add_card(slide, start_left + (card_width + gap) * 2, top, card_width, card_height, "3. Enterprise Features", [
        "Implement secure authentication (OAuth / JWT) for fraud analysts.",
        "Publish a robust RESTful API for seamless integration with web, mobile, or merchant endpoints."
    ], accent_color=ACCENT_BLUE)

def create_end_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    left = Inches(1.5)
    top = Inches(1.8)
    width = Inches(10.333)
    height = Inches(3.9)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ACCENT_BLUE
    card.line.width = Pt(3)
    
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.8)
    tf.margin_bottom = Inches(0.5)
    
    p_title = tf.paragraphs[0]
    p_title.text = "QUESTIONS & ANSWERS"
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_PRIMARY
    p_title.space_after = Pt(24)
    p_title.alignment = PP_ALIGN.CENTER
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Thank you for your time!"
    p_sub.font.name = 'Segoe UI'
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = TEXT_SECONDARY
    p_sub.space_after = Pt(20)
    p_sub.alignment = PP_ALIGN.CENTER
    
    p_contact = tf.add_paragraph()
    p_contact.text = "Credit Card Fraud Detection System • AI Presentation"
    p_contact.font.name = 'Segoe UI'
    p_contact.font.size = Pt(12)
    p_contact.font.color.rgb = ACCENT_BLUE
    p_contact.alignment = PP_ALIGN.CENTER

def main():
    prs = Presentation()
    
    # Set to modern Widescreen (16:9) aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    create_title_slide(prs)
    create_problem_slide(prs)
    create_solution_slide(prs)
    create_tech_stack_slide(prs)
    create_pipeline_slide(prs)
    create_architecture_slide(prs)
    create_ui_slide(prs)
    create_features_slide(prs)
    create_future_slide(prs)
    create_end_slide(prs)
    
    output_path = "Credit_Card_Fraud_Presentation.pptx"
    prs.save(output_path)
    print(f"Premium Presentation generated successfully at: {output_path}")

if __name__ == '__main__':
    main()
